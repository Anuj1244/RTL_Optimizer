from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation(output_file="RTL_Optimizer_Presentation.ppt"):
    prs = Presentation()

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "RTL Optimizer: AI-Powered Dynamic Power Reduction"
    subtitle.text = "Automating Sequential Enable Injection with LLMs\n[Your Name/Team]"

    # --- Slide 2: The Problem ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "The Problem: High Dynamic Power"
    content = slide.placeholders[1].text_frame
    content.text = "Challenges in Modern Chip Design:"
    p = content.add_paragraph()
    p.text = "• High power consumption due to clocking/switching."
    p.level = 1
    p = content.add_paragraph()
    p.text = "• Manual RTL refactoring is slow and error-prone."
    p.level = 1
    p = content.add_paragraph()
    p.text = "• Traditional tools are brittle with complex RTL styles."
    p.level = 1

    # --- Slide 3: The Solution ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "The Solution: AI-Powered RTL Optimizer"
    content = slide.placeholders[1].text_frame
    p = content.add_paragraph()
    p.text = "• Automated Workflow: From power reports to optimized RTL."
    p = content.add_paragraph()
    p.text = "• AI Integration: Llama 3.3 70B for context-aware refactoring."
    p = content.add_paragraph()
    p.text = "• Robustness: Handles diverse coding styles and provides validation."

    # --- Slide 4: System Workflow ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "System Architecture & Workflow"
    content = slide.placeholders[1].text_frame
    steps = ["1. Report Parsing", "2. Context Extraction", "3. LLM Refactoring", "4. Code Synthesis", "5. Validation & Diffing"]
    for step in steps:
        p = content.add_paragraph()
        p.text = step

    # --- Slide 5: Role of LLM ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Role of LLM (Llama 3.3 70B)"
    content = slide.placeholders[1].text_frame
    p = content.add_paragraph()
    p.text = "Why LLM?"
    p = content.add_paragraph()
    p.text = "• Context Awareness: Understands complex always-block structures."
    p.level = 1
    p = content.add_paragraph()
    p.text = "• Precision: Maintains original logic while adding enables."
    p.level = 1
    p = content.add_paragraph()
    p.text = "• Strategy: Specifically used for the 'Refactoring' step."
    p.level = 1

    # --- Slide 6: Step-by-Step I/O ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Step-by-Step Input & Output"
    content = slide.placeholders[1].text_frame
    io_data = [
        "Step 1: Input: Power Report (.txt) | Output: Metadata dict",
        "Step 2: Input: RTL + Metadata | Output: Always-block context",
        "Step 3 (LLM): Input: Context + Info | Output: Optimized RTL line",
        "Step 4: Input: RTL + LLM Patch | Output: Final optimized file"
    ]
    for item in io_data:
        p = content.add_paragraph()
        p.text = item
        p.font.size = Pt(18)

    # --- Slide 7: Results ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Validation & Results"
    content = slide.placeholders[1].text_frame
    p = content.add_paragraph()
    p.text = "• Automated Lint Checks: Balanced blocks and Verilog syntax."
    p = content.add_paragraph()
    p.text = "• Visual Verification: Git-style diffs for immediate review."
    p = content.add_paragraph()
    p.text = "• Efficiency: Drastically reduces manual engineering time."

    # --- Slide 8: Conclusion ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Conclusion"
    content = slide.placeholders[1].text_frame
    p = content.add_paragraph()
    p.text = "AI-Powered RTL Optimizer provides a scalable, reliable, and efficient way to reduce dynamic power in digital designs."
    p.font.bold = True

    prs.save(output_file)
    print(f"Presentation saved as {output_file}")

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("Error: 'python-pptx' library not found. Please install it using 'pip install python-pptx'.")
